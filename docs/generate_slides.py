import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Premium Dark Theme Palette (Sleek Slate-950 + Neon Emerald, Cyan, Purple & Amber)
    BG_DARK = RGBColor(11, 15, 25)           # Deep Space Slate
    CARD_DARK = RGBColor(22, 29, 45)         # Elevated Slate Card
    CARD_BORDER = RGBColor(45, 55, 75)       # Subtle Card Border
    CARD_HOVER = RGBColor(29, 39, 60)        # Card Header / Accent BG

    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_LIGHT = RGBColor(241, 245, 249)     # Slate 100
    TEXT_MUTED = RGBColor(148, 163, 184)     # Slate 400
    TEXT_SUBTLE = RGBColor(100, 116, 139)    # Slate 500

    EMERALD = RGBColor(52, 211, 153)         # Neon Emerald
    EMERALD_BG = RGBColor(6, 78, 59)         # Dark Emerald Container
    CYAN = RGBColor(56, 189, 248)            # Vibrant Sky/Cyan
    CYAN_BG = RGBColor(12, 74, 110)          # Dark Cyan Container
    PURPLE = RGBColor(192, 132, 252)         # Vibrant Purple
    PURPLE_BG = RGBColor(88, 28, 135)        # Dark Purple Container
    AMBER = RGBColor(251, 191, 36)           # Amber / Gold
    AMBER_BG = RGBColor(120, 53, 15)         # Dark Amber Container
    ROSE = RGBColor(251, 113, 133)           # Rose / Red Alert

    def add_slide_with_background(prs):
        slide = prs.slides.add_slide(blank_layout)
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.color.rgb = BG_DARK
        
        # Add smooth fade transition
        try:
            trans = parse_xml(f'<p:transition {nsdecls("p")} spd="med"><p:fade/></p:transition>')
            slide._element.append(trans)
        except Exception:
            pass
        return slide

    def add_header(slide, title_text, category_text="OPERATIONS & DYNAMIC SCHEDULING PLATFORM", slide_num=None):
        # Category Pill
        cat_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.32))
        cat_card.fill.solid()
        cat_card.fill.fore_color.rgb = CARD_HOVER
        cat_card.line.color.rgb = CARD_BORDER
        tf = cat_card.text_frame
        tf.margin_top = Inches(0.02)
        tf.margin_left = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = EMERALD

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(10.5), Inches(0.8))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

        # Slide Number Badge
        if slide_num:
            num_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.6), Inches(0.4), Inches(0.9), Inches(0.32))
            num_card.fill.solid()
            num_card.fill.fore_color.rgb = CARD_HOVER
            num_card.line.color.rgb = CARD_BORDER
            tf_num = num_card.text_frame
            tf_num.margin_top = Inches(0.02)
            num_p = tf_num.paragraphs[0]
            num_p.text = f"{slide_num} / 13"
            num_p.alignment = PP_ALIGN.CENTER
            num_p.font.size = Pt(10)
            num_p.font.bold = True
            num_p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = add_slide_with_background(prs)

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.9), Inches(5.2), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = CARD_HOVER
    badge.line.color.rgb = EMERALD
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "AIUB • SOFTWARE ENGINEERING (SECTION: DD • GROUP: H)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Agentic Operations & Dynamic\nScheduling Platform for Field Services"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    sub_box = s1.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Eliminating Double-Bookings, Customer No-Shows & Revenue Leakage via Deterministic State Machine Automation"
    p.font.size = Pt(14)
    p.font.color.rgb = CYAN

    # Left Card: Supervised by
    sup_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(4.5), Inches(2.0))
    sup_card.fill.solid()
    sup_card.fill.fore_color.rgb = CARD_DARK
    sup_card.line.color.rgb = CARD_BORDER
    tf = sup_card.text_frame
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "SUPERVISED BY"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p2 = tf.add_paragraph()
    p2.text = "Sourav Akib Sarkar"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(4)
    p3 = tf.add_paragraph()
    p3.text = "Faculty of Science & Technology\nDepartment of Computer Science, AIUB"
    p3.font.size = Pt(10.5)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(4)

    # Right Card: Development Team
    team_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(4.7), Inches(6.9), Inches(2.0))
    team_card.fill.solid()
    team_card.fill.fore_color.rgb = CARD_DARK
    team_card.line.color.rgb = CARD_BORDER
    tf = team_card.text_frame
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "DEVELOPMENT TEAM (GROUP H)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    members = [
        "1. Abrar Tajwar Khan (24-57356-2)",
        "2. Ahmed Aahan Addin (24-57325-2)",
        "3. Lasker Sifwat Hossein (24-57368-2)",
        "4. Sadnan Shahad (24-57524-2)",
    ]
    for m in members:
        p_m = tf.add_paragraph()
        p_m.text = m
        p_m.font.size = Pt(11)
        p_m.font.color.rgb = TEXT_LIGHT
        p_m.space_before = Pt(2)

    s1.notes_slide.notes_text_frame.text = "Good day faculty and peers. We present Group H's Software Engineering project: Agentic Operations and Dynamic Scheduling Platform for Field Services, developed under the supervision of Sourav Akib Sarkar."

    # =========================================================================
    # SLIDE 2: THE REAL-WORLD PROBLEM
    # =========================================================================
    s2 = add_slide_with_background(prs)
    add_header(s2, "The Real-World Operational Dilemma in Field Services", "PROBLEM STATEMENT", "2")

    problems = [
        ("01", "Double-Booking & Scheduling Collisions", "Manual phone calls & Excel sheets cause simultaneous overlapping customer bookings, unoptimized schedule gaps, and technician idle time.", ROSE),
        ("02", "High Customer No-Show Rates", "Without upfront financial commitment, customers cancel last-minute after technicians spend over an hour traveling through Dhaka traffic.", AMBER),
        ("03", "Blind Resource Dispatching", "Operations managers lack real-time visibility into technician skill specialties, active availability, and geographic zone buffers.", CYAN),
        ("04", "Manual Revenue Leakage", "Delayed paper billing, unrecorded cash balances, and lack of verifiable proof-of-work lead to disputed invoices and lost revenue.", PURPLE),
    ]

    for idx, (num, title, desc, accent) in enumerate(problems):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.8 + row * 2.5)

        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.65), Inches(2.25))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = f"PROBLEM {num}"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(13.5)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(3)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(6)

    s2.notes_slide.notes_text_frame.text = "In field operations, over 60% of time is wasted due to fragmented communication. Double bookings, customer no-shows after travel, blind dispatching, and manual revenue leakage are the 4 core problems we target."

    # =========================================================================
    # SLIDE 3: OUR UNIFIED ARCHITECTURE
    # =========================================================================
    s3 = add_slide_with_background(prs)
    add_header(s3, "System Solution: 3 Portals + 1 Autonomous Engine", "SYSTEM ARCHITECTURE", "3")

    portals = [
        ("Customer Portal", "/book • /track", "• Dynamic 30-min slot calendar\n• 15-minute travel buffers\n• 20% upfront deposit checkout (bKash/Nagad)\n• Real-time progress & SMS notification feed", EMERALD),
        ("Operations Dispatcher", "/admin/dispatch", "• 6-Column Kanban visual pipeline\n• Skill-based staff allocation\n• Dhaka zone proximity matching\n• Concurrency lock verification", PURPLE),
        ("Field Technician Portal", "/tech/active-job", "• 390px responsive mobile layout\n• Top numbered task switcher\n• Live on-site work duration stopwatch\n• HTML5 digital client signature sign-off", CYAN),
    ]

    for idx, (title, route, desc, accent) in enumerate(portals):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.8)

        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.75), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = route
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_LIGHT
        p3.space_before = Pt(12)

    s3.notes_slide.notes_text_frame.text = "Our platform connects three dedicated portals into a synchronized real-time system: Customer booking with deposits, Operations Dispatch Kanban, and Field Technician mobile sign-off."

    # =========================================================================
    # SLIDE 4: SYSTEM ACTORS & USE CASE MODEL
    # =========================================================================
    s4 = add_slide_with_background(prs)
    add_header(s4, "System Actors & Role Responsibilities (SRS Alignment)", "ACTORS & ROLES", "4")

    actors = [
        ("Customer (Tanvir Ahmed)", "• Initiates service requests (UC-01)\n• Pays upfront deposit online (US-02)\n• Tracks real-time status & ETA (UC-02)\n• Accesses historical receipts (UC-03)", "HUMAN ACTOR", CYAN),
        ("Field Technician (Kazi Shakil)", "• Receives assigned tasks on mobile\n• Advances linear states (UC-06)\n• Navigates & runs work timer\n• Collects digital sign-offs (NFR-04)", "HUMAN ACTOR", EMERALD),
        ("Admin / Dispatcher (Tajwar Hossain)", "• Matches staff by skill & zone (UC-05)\n• Manages Kanban pipeline (US-05)\n• Inspects tax invoices & revenue (UC-08)\n• Monitors productivity & CSAT (UC-09)", "HUMAN ACTOR", PURPLE),
        ("AI Scheduling Agent (System Worker)", "• Computes conflict-free slots (UC-04)\n• Enforces concurrency locks (FR-03)\n• Dispatches CRM webhook alerts (UC-07)\n• Auto-settles final tax invoices (FR-07)", "SYSTEM ACTOR (4th)", AMBER),
    ]

    for idx, (title, desc, badge_txt, accent) in enumerate(actors):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.8 + row * 2.5)

        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.65), Inches(2.25))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.22)
        tf.margin_right = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = badge_txt
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(13.5)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(3)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = TEXT_LIGHT
        p3.space_before = Pt(5)

    s4.notes_slide.notes_text_frame.text = "In our SRS, we define 3 human login personas and 1 automated system actor: the AI Scheduling Agent that autonomously evaluates calendar states, enforces race condition locks, and settles invoices."

    # =========================================================================
    # SLIDE 5: METHODOLOGY — AGILE SCRUM
    # =========================================================================
    s5 = add_slide_with_background(prs)
    add_header(s5, "Software Engineering Methodology: Agile Scrum", "DEVELOPMENT PROCESS", "5")

    scrum_pillars = [
        ("01", "Why Scrum?", "Iterative delivery across 4 sprints, continuous requirement traceability against SRS-2, and frequent UI/UX validation.", EMERALD),
        ("02", "Product Backlog", "Structured user stories directly mapped to functional requirements FR-01 through FR-08 with clear acceptance criteria.", CYAN),
        ("03", "Sprint Ceremonies", "Conducted bi-weekly Sprint Planning, daily peer standups, sprint reviews, and retrospective code audits.", PURPLE),
        ("04", "Quality & Burndown", "Monitored velocity, enforced zero build/lint errors (npm run build), and validated 100% route health.", AMBER),
    ]

    for idx, (num, title, desc, accent) in enumerate(scrum_pillars):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)

        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.75), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_LIGHT
        p3.space_before = Pt(12)

    s5.notes_slide.notes_text_frame.text = "We adopted the Scrum framework to build this system iteratively. Product backlog items were derived directly from SRS requirements and delivered in 4 two-week sprint increments."

    # =========================================================================
    # SLIDE 6: 4-SPRINT LIFECYCLE BREAKDOWN
    # =========================================================================
    s6 = add_slide_with_background(prs)
    add_header(s6, "4-Sprint Incremental Lifecycle Breakdown", "SCRUM EXECUTION", "6")

    sprints = [
        ("Sprint 1 (Weeks 1-2)", "Architecture & Database Foundations", "• SRS requirements analysis & schema DDL\n• RoleGuard RBAC & unified login portal\n• Next.js 14 App Router + Tailwind setup\n• PostgreSQL seed data & state schemas", EMERALD),
        ("Sprint 2 (Weeks 3-4)", "Dynamic Slot Engine & Booking", "• Pure TypeScript slot algorithm (scheduling.ts)\n• 30-min window evaluation & 15-min buffers\n• 3-Step booking wizard with 20% deposit\n• Concurrency locking badge (FR-03)", CYAN),
        ("Sprint 3 (Weeks 5-6)", "Finite State Machine & Mobile", "• Strict linear FSM engine (state-machine.ts)\n• 6-Column Operations Dispatch Kanban\n• Skill-based staff assignment & Dhaka zones\n• 390px Mobile portal + Signature pad", PURPLE),
        ("Sprint 4 (Weeks 7-8)", "CRM Webhooks & Analytics", "• Automated CRM SMS/WhatsApp alert stream\n• Autonomous tax invoicing modal (INV-...)\n• Revenue & Technician Productivity metrics\n• Production build verification (npm run build)", AMBER),
    ]

    for idx, (sprint_name, sprint_goal, sprint_items, accent) in enumerate(sprints):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)

        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.75), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = sprint_name
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = sprint_goal
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = sprint_items
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = TEXT_LIGHT
        p3.space_before = Pt(8)

    s6.notes_slide.notes_text_frame.text = "Across Sprint 1 through 4, we transitioned from database architecture to slot calculations, then to the dispatch Kanban and mobile portal, and finally to automated CRM alerts, invoicing, and analytics."

    # =========================================================================
    # SLIDE 7: FEATURE 1 — DYNAMIC SLOT ENGINE & LOCKING
    # =========================================================================
    s7 = add_slide_with_background(prs)
    add_header(s7, "Feature 1: Intelligent Dynamic Slot Engine & Concurrency Locking", "CORE FEATURE", "7")

    card_left = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.65), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_DARK
    card_left.line.color.rgb = CARD_BORDER
    tf = card_left.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "Algorithmic Slot Calculation (FR-02)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p2 = tf.add_paragraph()
    p2.text = "• Operates over business hours: 8:00 AM – 6:00 PM\n• Generates 30-minute start time candidates\n• Evaluates exact service duration (45 to 180 mins)\n• Appends mandatory 15-minute travel & prep buffer\n• Filters past slots dynamically for same-day bookings\n• Matches technician capacity in real-time"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)

    card_right = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_DARK
    card_right.line.color.rgb = CARD_BORDER
    tf2 = card_right.text_frame
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)
    p = tf2.paragraphs[0]
    p.text = "Concurrency Locking & Deposit (FR-03, FR-06)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p2 = tf2.add_paragraph()
    p2.text = "• Overlap Condition: (Start < BookEnd) AND (End > BookStart)\n• Database transactional locking eliminates double-booking collisions\n• 20% Upfront Deposit Gateway (bKash / Nagad / Card)\n• Eliminates customer no-shows through financial commitment\n• Creates immutable booking record in 'Pending' state"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)

    s7.notes_slide.notes_text_frame.text = "The slot engine evaluates operating hours and service duration with a 15-minute buffer. Transactional locks prevent double bookings, and the 20% deposit secures customer commitment."

    # =========================================================================
    # SLIDE 8: FEATURE 2 — FINITE STATE MACHINE
    # =========================================================================
    s8 = add_slide_with_background(prs)
    add_header(s8, "Feature 2: Deterministic Finite State Machine Lifecycle", "CORE FEATURE", "8")

    fsm_steps = [
        ("1. Pending", "Deposit locked; awaiting dispatcher assign", "Customer", CYAN),
        ("2. Scheduled", "Staff assigned; slot & zone locked", "Dispatcher", PURPLE),
        ("3. En Route", "Technician departed; arrival SMS sent", "Technician", AMBER),
        ("4. In Progress", "Tech on-site; work timer running", "Technician", EMERALD),
        ("5. Completed", "Job finished; digital signature captured", "Technician", CYAN),
        ("6. Billed", "Automated invoice generated & settled", "System Bot", EMERALD),
    ]

    for idx, (title, desc, actor, accent) in enumerate(fsm_steps):
        left = Inches(0.8 + idx * 1.95)
        top = Inches(1.8)

        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.2)
        tf.margin_right = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = f"STEP {idx+1}"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = f"Actor: {actor}"
        p3.font.size = Pt(9.5)
        p3.font.bold = True
        p3.font.color.rgb = CYAN
        p3.space_before = Pt(8)

        p4 = tf.add_paragraph()
        p4.text = desc
        p4.font.size = Pt(9.5)
        p4.font.color.rgb = TEXT_MUTED
        p4.space_before = Pt(8)

    s8.notes_slide.notes_text_frame.text = "Our state engine enforces strict linear progression: Pending, Scheduled, En Route, In Progress, Completed, Billed. Illegal state jumps are strictly rejected by the validator."

    # =========================================================================
    # SLIDE 9: FEATURE 3 — SKILL DISPATCH & MOBILE FIELD
    # =========================================================================
    s9 = add_slide_with_background(prs)
    add_header(s9, "Feature 3: Skill-Based Dispatch & Technician Mobile View", "CORE FEATURE", "9")

    card1 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.65), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_DARK
    card1.line.color.rgb = CARD_BORDER
    tf = card1.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "Skill-Based Dispatching (UC-05)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p2 = tf.add_paragraph()
    p2.text = "• 6-Column Kanban Board representing the FSM pipeline\n• Technician skill taxonomy (Paint Buffing, Steam Sanitization, Diagnostics)\n• Automatic '✓ Skill Match' badge on assignment modal\n• Geographical zone proximity routing (Gulshan / Banani / Dhanmondi)\n• Real-time technician active workload indicator"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)

    card2 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_DARK
    card2.line.color.rgb = CARD_BORDER
    tf2 = card2.text_frame
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)
    p = tf2.paragraphs[0]
    p.text = "Field Technician Mobile Portal (NFR-04)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p2 = tf2.add_paragraph()
    p2.text = "• 390px mobile viewport frame with minimal wireframe aesthetics\n• Top Numbered Task Switcher Bar (Task 1, Task 2, Task 3)\n• 1-Tap Sequential Actions: 'Accept & Start Trip' ➔ 'Arrived & Start Timer'\n• Live on-site work duration stopwatch\n• HTML5 signature canvas for digital client sign-off"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)

    s9.notes_slide.notes_text_frame.text = "Dispatchers assign technicians matching service skills and location zones. Technicians execute jobs on their mobile portal with task switching, on-site timer, and client signature sign-off."

    # =========================================================================
    # SLIDE 10: FEATURE 4 — CRM ALERTS, INVOICING & ANALYTICS
    # =========================================================================
    s10 = add_slide_with_background(prs)
    add_header(s10, "Feature 4: Automated CRM Alerts, Tax Invoices & Analytics", "CORE FEATURE", "10")

    card1 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.65), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_DARK
    card1.line.color.rgb = CARD_BORDER
    tf = card1.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "Automated CRM Webhook Alerts (FR-05)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p2 = tf.add_paragraph()
    p2.text = "• Event-driven SMS/WhatsApp alert dispatch to customer phone\n• Alert 1: Deposit Confirmed (Awaiting Dispatcher)\n• Alert 2: Staff Scheduled (Technician Assigned)\n• Alert 3: En Route ETA Update (~20 mins)\n• Alert 4: Completed & Official Tax Invoice Ready"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)

    card2 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_DARK
    card2.line.color.rgb = CARD_BORDER
    tf2 = card2.text_frame
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)
    p = tf2.paragraphs[0]
    p.text = "Tax Invoices & Analytics (FR-07, FR-08)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = AMBER
    p2 = tf2.add_paragraph()
    p2.text = "• Autonomous billing settlement upon signature capture\n• Itemized Printable Tax Invoice Modal (INV-...)\n• Real-Time Revenue Pipeline in BDT (৳)\n• Technician Productivity Table (Utilization %, Active vs Completed)\n• Customer Feedback CSAT Rating (4.9 / 5.0 ★ • 98% Satisfaction)"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)

    s10.notes_slide.notes_text_frame.text = "The CRM webhook engine sends live SMS updates. Upon signature capture, tax invoices are auto-generated and operational analytics track gross pipeline revenue and technician utilization."

    # =========================================================================
    # SLIDE 11: LIVE DEMO WALKTHROUGH
    # =========================================================================
    s11 = add_slide_with_background(prs)
    add_header(s11, "Live Demonstration Walkthrough (5-Step Pipeline)", "SYSTEM DEMO", "11")

    demo_steps = [
        ("Step 1", "Clean Zero-Job Pipeline", "Open /admin/dispatch to verify clean 0-job state ready for demo.", CYAN),
        ("Step 2", "Customer Booking (/book)", "Select service, pick dynamic slot, pay 20% deposit via bKash/Nagad.", EMERALD),
        ("Step 3", "Dispatcher Assignment (/admin/dispatch)", "Assign technician matching skill taxonomy & zone. Moves to Scheduled.", PURPLE),
        ("Step 4", "Mobile Execution (/tech/active-job)", "Accept job, start trip (En Route), run work timer (In Progress), collect digital signature.", AMBER),
        ("Step 5", "Automated Invoicing & Analytics (/admin/analytics)", "Auto-settles to Billed, generates printable invoice, and updates revenue.", EMERALD),
    ]

    for idx, (step_num, title, desc, accent) in enumerate(demo_steps):
        top = Inches(1.8 + idx * 0.95)
        card = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(0.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = f"{step_num}: {title} — "
        p.font.size = Pt(11)
        p.font.color.rgb = accent
        p.font.bold = True
        
        # description
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_LIGHT

    s11.notes_slide.notes_text_frame.text = "Now we will switch to our live browser window at localhost:3000 to demonstrate the full 5-step lifecycle."

    # =========================================================================
    # SLIDE 12: TECH STACK & PRODUCTION VERIFICATION
    # =========================================================================
    s12 = add_slide_with_background(prs)
    add_header(s12, "Technology Stack & Verification Metrics", "TECH STACK", "12")

    tech_cards = [
        ("Frontend & UI", "• Next.js 14 (App Router)\n• React 18 & TypeScript\n• TailwindCSS & Lucide Icons\n• Responsive 390px Mobile View", CYAN),
        ("Backend & Logic", "• PostgreSQL / Supabase Schema\n• Finite State Machine Engine\n• Dynamic Slot Algorithm\n• Webhook REST Endpoints", PURPLE),
        ("State Management", "• Reactive Context Store (store.tsx)\n• LocalStorage persistence\n• Zero third-party runtime failure\n• Offline & Vercel ready", AMBER),
        ("Production Verification", "• npm run build: 12/12 routes\n• 0 Build & 0 Type Errors\n• 100% SRS-2 Compliance\n• Pushed to GitHub (main)", EMERALD),
    ]

    for idx, (title, desc, accent) in enumerate(tech_cards):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)

        card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.75), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_LIGHT
        p2.space_before = Pt(12)

    s12.notes_slide.notes_text_frame.text = "Our platform is built with Next.js 14, TypeScript, and TailwindCSS. The production build passed with zero errors across all 12 routes."

    # =========================================================================
    # SLIDE 13: CONCLUSION & Q&A
    # =========================================================================
    s13 = add_slide_with_background(prs)

    badge = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(3.8), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = CARD_HOVER
    badge.line.color.rgb = EMERALD
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJECT SUMMARY & CONCLUSION"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    title_box = s13.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!\nQuestions & Discussion"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    card = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.8), Inches(11.7), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_DARK
    card.line.color.rgb = CARD_BORDER
    tf = card.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "KEY TAKEAWAYS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    bullets = [
        "✔ 100% compliance with SRS-2 scopes, use cases (UC-01 to UC-09), and functional requirements (FR-01 to FR-08).",
        "✔ Successfully eliminates double-bookings, customer no-shows, and revenue leakage.",
        "✔ Developed iteratively following Agile Scrum across 4 two-week sprint milestones.",
        "✔ Fully deployable on Vercel with zero external database setup barriers.",
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(11.5)
        p_b.font.color.rgb = TEXT_LIGHT
        p_b.space_before = Pt(4)

    s13.notes_slide.notes_text_frame.text = "In conclusion, our platform successfully achieves all goals set out in our SRS. Thank you, and we now welcome any questions from the honorable faculty."

    output_path = "/Users/tajwar/Desktop/SWE Project/docs/presentation_slides.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
